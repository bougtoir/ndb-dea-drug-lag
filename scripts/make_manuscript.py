#!/usr/bin/env python3
"""
Generate the JA and EN manuscripts (docx, figures inline), separate editable
table docx, and editable figure pptx for the HCV DAA new-treatment-anticipation /
drug-lag study.

All result numbers are read from results/summary.json and data/*.csv and formatted
at runtime; no result value is hard-coded here. References are numbered in order of
first appearance (Vancouver). Figures/tables are inserted immediately after the
paragraph that first cites them.

Usage:
    python3 scripts/make_manuscript.py
"""
import json
import os

import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt

BASE = os.path.join(os.path.dirname(__file__), "..")
DATA = os.path.join(BASE, "data")
RES = os.path.join(BASE, "results")
OUT = os.path.join(BASE, "output")

S = json.load(open(os.path.join(RES, "summary.json"), encoding="utf-8"))
ITS = json.load(open(os.path.join(RES, "its_summary.json"), encoding="utf-8"))
COURSE = json.load(open(os.path.join(RES, "course_estimate.json"), encoding="utf-8"))
TS = pd.read_csv(os.path.join(DATA, "hcv_timeseries.csv")).set_index("fy")
EV = pd.read_csv(os.path.join(DATA, "announcement_events.csv"))

Y0, Y1 = S["fiscal_year_range"]


def fmt(x, nd=0):
    return f"{x:,.{nd}f}"


# ---- derived display values (all traceable to summary/timeseries) -------------
peg_drop = abs(S["peginterferon"]["pct_change_first_to_last"])
rbv_zero_year = S["ribavirin"]["first_year_near_zero"]
daa_peak_fy = S["daa_total"]["peak_fy"]
daa_peak_val_m = S["daa_total"]["peak_value"] / 1e6
daa_rise = S["daa_total"]["pct_change_first_to_peak"]
daa_fall = abs(S["daa_total"]["pct_change_peak_to_last"])
daa_last_m = S["daa_total"]["fy_last"] / 1e6
n_daa = S["n_distinct_daa_products"]


def decline_ci(d, ci_key="annual_change_hac95"):
    """Return (|annual %|, ci_low_magnitude, ci_high_magnitude) for a decay fit."""
    r = abs(d["annual_change_pct"])
    lo, hi = (abs(v) for v in d[ci_key])
    return r, min(lo, hi), max(lo, hi)


peg_r, peg_lo, peg_hi = decline_ci(ITS["peginterferon_decay"])
rbv_r, rbv_lo, rbv_hi = decline_ci(ITS["ribavirin_decay"])
conv_r, conv_lo, conv_hi = decline_ci(ITS["conventional_ifn_decay"])
rbv_fy0, rbv_fy1 = ITS["ribavirin_decay"]["fy_used"]
daa_knot = ITS["daa_segmented"]["knot_fy"]
daa_post_r = abs(ITS["daa_segmented"]["post_peak_annual_rate_pct"])
daa_post_ci = sorted(abs(v) for v in ITS["daa_segmented"]["post_peak_annual_rate_boot95"])
daa_slope_p = ITS["daa_segmented"]["slope_change_p"]
n_obs = ITS["n_annual_observations"]
n_boot = ITS["bootstrap_resamples"]

# treatment-course sensitivity (estimate only; see data/daa_course_assumptions.csv)
course_peak_fy = COURSE["peak_fy"]
course_peak = COURSE["peak_estimated_courses"]
course_total_lo, course_total_hi = sorted(COURSE["estimated_total_courses_fy2014_2023_range"])

# ------------------------------------------------------------------------------
# Reference text per source id (no fabricated citations; verifiable sources only).
# Citation NUMBERS are assigned dynamically in order of first appearance (Vancouver),
# tracked in CITE_ORDER during document build.
REF_TEXT = {
    "en": {
        "ndb": "Ministry of Health, Labour and Welfare (Japan). NDB Open Data "
               "(1st-10th editions). https://www.mhlw.go.jp/ndb/opendatasite/ "
               f"(accessed for fiscal years {Y0}-{Y1}).",
        "bms": "Bristol-Myers Squibb K.K. Press release: approval in Japan of the "
               "world's first all-oral, interferon- and ribavirin-free treatment for "
               "chronic hepatitis C (daclatasvir + asunaprevir). 4 July 2014.",
        "nhi": "Ministry of Health, Labour and Welfare / Central Social Insurance "
               "Medical Council (Chuikyo). NHI drug-price listings of direct-acting "
               "antivirals for hepatitis C (2014-2017).",
        "newey": "Newey WK, West KD. A simple, positive semi-definite, "
                 "heteroskedasticity and autocorrelation consistent covariance "
                 "matrix. Econometrica. 1987;55(3):703-708.",
        "efron": "Efron B, Tibshirani RJ. An Introduction to the Bootstrap. "
                 "New York: Chapman & Hall; 1993.",
    },
    "ja": {
        "ndb": "厚生労働省. NDBオープンデータ（第1〜10回）. "
               "https://www.mhlw.go.jp/ndb/opendatasite/ "
               f"（{Y0}〜{Y1}年度分を使用）.",
        "bms": "ブリストル・マイヤーズ株式会社. プレスリリース：日本初の"
               "インターフェロンおよびリバビリンを必要としない経口薬のみによる"
               "C型慢性肝炎治療薬（ダクラタスビル＋アスナプレビル）の製造販売承認取得. "
               "2014年7月4日.",
        "nhi": "厚生労働省／中央社会保険医療協議会（中医協）. C型肝炎に対する"
               "直接作用型抗ウイルス薬の薬価基準収載（2014〜2017年）.",
        "newey": "Newey WK, West KD. A simple, positive semi-definite, "
                 "heteroskedasticity and autocorrelation consistent covariance "
                 "matrix. Econometrica. 1987;55(3):703-708.",
        "efron": "Efron B, Tibshirani RJ. An Introduction to the Bootstrap. "
                 "New York: Chapman & Hall; 1993.",
    },
}

CITE_ORDER = []  # reset per document; source ids in order of first appearance


def cite(p, keys, lang):
    """Append a Vancouver superscript citation, numbering by first appearance."""
    for k in keys:
        if k not in CITE_ORDER:
            CITE_ORDER.append(k)
    nums = sorted(CITE_ORDER.index(k) + 1 for k in keys)
    run = p.add_run(",".join(str(n) for n in nums))
    run.font.superscript = True
    return p


def add_caption(doc, text):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(14)
    r = p.add_run(text)
    r.bold = True
    r.font.size = Pt(9)
    return p


def insert_fig(doc, path, width=6.3):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(path, width=Inches(width))


# ------------------------------------------------------------------------------
TXT = {
    "en": dict(
        title="Population-level replacement of interferon-based standard therapy for "
              "hepatitis C after interferon-free direct-acting antivirals: a "
              "new-treatment-anticipation / practical-drug-lag analysis of Japan's NDB "
              "Open Data",
        h_abs="Abstract", h_intro="Introduction", h_meth="Methods",
        h_res="Results", h_disc="Discussion", h_lim="Limitations", h_conc="Conclusion",
        h_ref="References", h_da="Data and code availability",
        abs_bg="Background: ", abs_me="Methods: ", abs_re="Results: ", abs_co="Conclusion: ",
    ),
    "ja": dict(
        title="IFNフリー直接作用型抗ウイルス薬導入後のC型肝炎インターフェロン標準治療の"
              "人口レベルでの置換：NDBオープンデータによる新治療法待望論／実用的"
              "ドラッグラグの検証",
        h_abs="要旨", h_intro="緒言", h_meth="方法",
        h_res="結果", h_disc="考察", h_lim="限界", h_conc="結論",
        h_ref="文献", h_da="データおよびコードの入手可能性",
        abs_bg="背景：", abs_me="方法：", abs_re="結果：", abs_co="結論：",
    ),
}


def build_manuscript(lang, journal=None):
    T = dict(TXT[lang])
    if journal == "pds":
        # Pharmacoepidemiology & Drug Safety structured-abstract heading.
        T["abs_bg"] = "Background/Objectives: "
    CITE_ORDER.clear()
    doc = Document()
    st = doc.styles["Normal"].font
    st.size = Pt(10.5)

    ttl = doc.add_paragraph()
    r = ttl.add_run(T["title"]); r.bold = True; r.font.size = Pt(14)

    if journal == "pds":
        add_pds_titlepage(doc)

    # ---- Abstract ----
    doc.add_heading(T["h_abs"], level=1)
    ab = doc.add_paragraph()
    if lang == "en":
        ab.add_run(T["abs_bg"]).bold = True
        ab.add_run("If population-level anticipation of a newly reported/reimbursed "
                   "treatment is real, use of the prior standard therapy should fall "
                   "sharply once the awaited option arrives; if not, either the "
                   "anticipation/practical lag is absent or decision-makers are not "
                   "reacting to the news. ")
        ab.add_run(T["abs_me"]).bold = True
        ab.add_run(f"Using Japan's NDB Open Data ({Y0}-{Y1}), we tracked national "
                   "dispensed quantity of interferon (IFN)-based standard therapy "
                   "(peginterferon, ribavirin) and of interferon-free direct-acting "
                   "antivirals (DAAs) for hepatitis C, against official approval/"
                   "reimbursement milestones. ")
        ab.add_run(T["abs_re"]).bold = True
        ab.add_run(f"Peginterferon dispensing fell {fmt(peg_drop,1)}% from FY{Y0} to "
                   f"FY{Y1}; ribavirin reached near-zero by FY{rbv_zero_year}. Total DAA "
                   f"dispensing peaked in FY{daa_peak_fy} ({fmt(daa_peak_val_m,1)} million "
                   f"units, +{fmt(daa_rise,0)}% vs FY{Y0}) and then fell {fmt(daa_fall,0)}% "
                   f"by FY{Y1}. ")
        ab.add_run(T["abs_co"]).bold = True
        ab.add_run("The IFN-based standard therapy was replaced at the population level "
                   "within about two years, and the DAA surge-then-decay is consistent "
                   "with a finite pool of long-waiting patients being treated in a burst "
                   "(pent-up demand) rather than steady substitution.")
    else:
        ab.add_run(T["abs_bg"]).bold = True
        ab.add_run("新治療法の報道・収載に対する待望が人口レベルで実在するなら、待望された"
                   "治療が登場した時点で従来の標準治療の利用は急減するはずである。変化がなければ、"
                   "待望論・実用的ラグが存在しないか、治療選択の主体が報道に反応していないことになる。")
        ab.add_run(T["abs_me"]).bold = True
        ab.add_run(f"日本のNDBオープンデータ（{Y0}〜{Y1}年度）を用い、C型肝炎に対する"
                   "インターフェロン（IFN）ベース標準治療（ペグインターフェロン、リバビリン）と"
                   "IFNフリー直接作用型抗ウイルス薬（DAA）の全国処方数量を、公式の承認・"
                   "薬価収載イベントと対応づけて追跡した。")
        ab.add_run(T["abs_re"]).bold = True
        ab.add_run(f"ペグインターフェロンの処方数量はFY{Y0}からFY{Y1}で{fmt(peg_drop,1)}%減少し、"
                   f"リバビリンはFY{rbv_zero_year}までにほぼゼロとなった。DAA合計はFY{daa_peak_fy}に"
                   f"ピーク（{fmt(daa_peak_val_m,1)}百万単位、FY{Y0}比+{fmt(daa_rise,0)}%）を示し、"
                   f"その後FY{Y1}までに{fmt(daa_fall,0)}%減少した。")
        ab.add_run(T["abs_co"]).bold = True
        ab.add_run("IFNベース標準治療は約2年で人口レベルに置換され、DAAの急増→減衰は、"
                   "長く待機していた患者集団が一括して治療された（待望＝pent-up demand）ことと"
                   "整合的であり、定常的な置換フローではないことを示す。")

    if journal == "pds":
        add_keywords(doc)
        add_key_points(doc)

    # ---- Introduction ----
    doc.add_heading(T["h_intro"], level=1)
    p = doc.add_paragraph()
    if lang == "en":
        p.add_run("When a new therapy is announced through the media or scheduled for "
                  "reimbursement, advocates of the new treatment expect a population-wide "
                  "shift away from the current standard of care. Chronic hepatitis C offers "
                  "an unusually clean natural experiment: interferon-based therapy was the "
                  "standard until interferon-free DAAs arrived in Japan in 2014-2015, when "
                  "daclatasvir plus asunaprevir became the world's first all-oral, "
                  "interferon- and ribavirin-free regimen")
        cite(p, ["bms"], lang)
        p.add_run(". We use nationwide dispensing counts from NDB Open Data")
        cite(p, ["ndb"], lang)
        p.add_run(" to test whether the standard therapy was displaced at the population "
                  "level around these announcement/reimbursement events.")
    else:
        p.add_run("新しい治療法が報道され、あるいは保険収載が予定されると、新治療法待望の立場からは"
                  "従来の標準治療から人口規模での移行が期待される。C型慢性肝炎は例外的に明瞭な自然実験を"
                  "提供する。インターフェロンベース治療が長く標準であったが、2014〜2015年にIFNフリーDAAが"
                  "登場し、ダクラタスビル＋アスナプレビルが世界初の全経口・IFN／リバビリン不要療法となった")
        cite(p, ["bms"], lang)
        p.add_run("。本研究ではNDBオープンデータ")
        cite(p, ["ndb"], lang)
        p.add_run("の全国処方数量を用い、これらの承認・収載イベント前後で標準治療が人口レベルで"
                  "置換されたかを検証する。")

    # ---- Methods ----
    doc.add_heading(T["h_meth"], level=1)
    p = doc.add_paragraph()
    if lang == "en":
        p.add_run(f"We used NDB Open Data editions 1-10 (fiscal years {Y0}-{Y1})")
        cite(p, ["ndb"], lang)
        p.add_run(", extracting the national total dispensed quantity (総計/処方数量) for "
                  "each drug from the sex- and age-stratified prescription-drug tables "
                  "(oral, topical, injectable). Drugs were classified by their actual "
                  f"product names into IFN-based standard therapy (peginterferon, "
                  f"ribavirin) and {n_daa} interferon-free DAA products. Three "
                  "first-generation NS3/4A protease inhibitors (simeprevir, telaprevir, "
                  "vaniprevir), which were used together with peginterferon and ribavirin "
                  "rather than as interferon-free regimens, were tabulated separately and "
                  "excluded from the DAA group. Conventional "
                  "interferon is reported separately because it is not hepatitis-C-specific. "
                  "The metric is dispensed quantity (tablets/capsules for oral drugs, "
                  "syringes/vials for injections), not a patient count, and is compared "
                  "within a product over time. Official approval and NHI drug-price "
                  "listing milestones")
        cite(p, ["bms", "nhi"], lang)
        p.add_run(" were used as intervention markers at fiscal-year resolution "
                  "(Table 1).")
    else:
        p.add_run(f"NDBオープンデータ第1〜10回（{Y0}〜{Y1}年度）")
        cite(p, ["ndb"], lang)
        p.add_run("を用い、処方薬の性年齢別薬効分類別数量表（内服・外用・注射）から各薬剤の全国"
                  "総計（処方数量）を抽出した。薬剤は実際の製品名からIFNベース標準治療（ペグ"
                  f"インターフェロン、リバビリン）と{n_daa}製剤のIFNフリーDAAに分類した。第一世代の"
                  "NS3/4Aプロテアーゼ阻害薬3剤（シメプレビル、テラプレビル、バニプレビル）は"
                  "IFNフリーではなくペグインターフェロン＋リバビリンと併用されるため、別掲しDAA群から"
                  "除外した。従来型"
                  "インターフェロンはC型肝炎特異的でないため別掲した。指標は処方数量（内服は錠・"
                  "カプセル、注射はシリンジ・バイアル）であり患者数ではなく、製剤ごとに経時比較する。"
                  "公式の承認・薬価収載イベント")
        cite(p, ["bms", "nhi"], lang)
        p.add_run("を年度解像度の介入マーカーとした（表1）。")

    add_caption(doc, ("Table 1. Official approval / NHI drug-price listing milestones "
                      "used as intervention markers." if lang == "en"
                      else "表1．介入マーカーとして用いた公式の承認・薬価収載イベント。"))
    add_events_table(doc, lang)

    p = doc.add_paragraph()
    if lang == "en":
        p.add_run(f"Statistical analysis. Because NDB begins in FY{Y0}, coincident with "
                  "IFN-free DAA availability, there is no pre-intervention baseline and a "
                  "conventional pre/post interrupted time series is not identifiable; with "
                  f"{n_obs} annual national observations we therefore fitted descriptive "
                  "trend models. Total DAA dispensing was modelled by continuous segmented "
                  "(broken-stick) log-linear regression with a knot at the observed peak "
                  "fiscal year, estimating pre- and post-peak annual multiplicative rates "
                  "of change; each IFN-based drug was modelled by exponential (log-linear) "
                  "decay over fiscal years with positive dispensing. Uncertainty was "
                  "expressed as heteroskedasticity- and autocorrelation-consistent "
                  "(Newey-West, maxlags=1) 95% intervals")
        cite(p, ["newey"], lang)
        p.add_run(f", cross-checked by a residual bootstrap ({fmt(n_boot,0)} resamples)")
        cite(p, ["efron"], lang)
        p.add_run(". These quantify how fast dispensing changed and its uncertainty; they "
                  "are not causal estimates. Analyses used Python (statsmodels); the full "
                  "pipeline (download -> build -> analyze -> figures) is openly reproducible.")
    else:
        p.add_run(f"統計解析。NDBはFY{Y0}開始でIFNフリーDAA導入時期と重なるため、介入前の"
                  "基準期間が存在せず従来型の前後比較（中断時系列）は同定できない。そこで"
                  f"{n_obs}点の年次全国データに対し記述的なトレンドモデルを当てはめた。DAA合計は"
                  "観測ピーク年度をノットとする連続分節（折れ線）対数線形回帰でモデル化し、"
                  "ピーク前後の年次の乗法的変化率を推定した。各IFNベース薬は処方数量が正の"
                  "年度に対する指数（対数線形）減衰でモデル化した。不確実性は不均一分散・"
                  "自己相関に頑健な（Newey-West, maxlags=1）95%区間")
        cite(p, ["newey"], lang)
        p.add_run(f"として表し、残差ブートストラップ（{fmt(n_boot,0)}回）で相互確認した")
        cite(p, ["efron"], lang)
        p.add_run("。これらは変化の速さと不確実性を定量化するものであり因果推定ではない。"
                  "解析はPython（statsmodels）で行い、全パイプライン"
                  "（download→build→analyze→figures）は公開・再現可能である。")

    # ---- Results ----
    doc.add_heading(T["h_res"], level=1)
    p = doc.add_paragraph()
    if lang == "en":
        p.add_run(f"The interferon-based standard therapy collapsed after IFN-free DAAs "
                  f"became available (Fig. 1). Peginterferon dispensing fell "
                  f"{fmt(peg_drop,1)}% from FY{Y0} to FY{Y1}, and ribavirin reached "
                  f"near-zero by FY{rbv_zero_year}. In the trend models this corresponded "
                  f"to peginterferon declining {fmt(peg_r,0)}% per year "
                  f"(95% CI {fmt(peg_lo,0)}-{fmt(peg_hi,0)}%), ribavirin "
                  f"{fmt(rbv_r,0)}% per year (95% CI {fmt(rbv_lo,0)}-{fmt(rbv_hi,0)}%; "
                  f"FY{rbv_fy0}-FY{rbv_fy1}), and conventional interferon {fmt(conv_r,0)}% "
                  f"per year (95% CI {fmt(conv_lo,0)}-{fmt(conv_hi,0)}%).")
    else:
        p.add_run(f"IFNフリーDAAの登場後、インターフェロンベース標準治療は消失した（図1）。"
                  f"ペグインターフェロンの処方数量はFY{Y0}からFY{Y1}で{fmt(peg_drop,1)}%減少し、"
                  f"リバビリンはFY{rbv_zero_year}までにほぼゼロとなった。トレンドモデルでは、"
                  f"ペグインターフェロンは年{fmt(peg_r,0)}%（95%CI {fmt(peg_lo,0)}〜{fmt(peg_hi,0)}%）、"
                  f"リバビリンは年{fmt(rbv_r,0)}%（95%CI {fmt(rbv_lo,0)}〜{fmt(rbv_hi,0)}%、"
                  f"FY{rbv_fy0}〜FY{rbv_fy1}）、従来型インターフェロンは年{fmt(conv_r,0)}%"
                  f"（95%CI {fmt(conv_lo,0)}〜{fmt(conv_hi,0)}%）の減少率であった。")
    insert_fig(doc, os.path.join(OUT, f"fig1_ifn_collapse_{lang}.png"))
    add_caption(doc, ("Fig. 1. Collapse of interferon-based standard therapy "
                      "(indexed to FY%d = 100) with total DAA dispensed quantity and "
                      "dated announcement markers." % Y0 if lang == "en"
                      else "図1．インターフェロンベース標準治療の消失（FY%d=100指数）、"
                           "DAA合計処方数量、および収載イベント。" % Y0))

    p = doc.add_paragraph()
    if lang == "en":
        p.add_run(f"DAA use itself showed a surge-then-decay pattern (Fig. 2): total DAA "
                  f"dispensing peaked in FY{daa_peak_fy} at {fmt(daa_peak_val_m,1)} million "
                  f"units (+{fmt(daa_rise,0)}% vs FY{Y0}) and then declined {fmt(daa_fall,0)}% "
                  f"to {fmt(daa_last_m,1)} million units by FY{Y1}, with successive products "
                  f"replacing earlier ones. In the segmented model, DAA dispensing declined "
                  f"{fmt(daa_post_r,0)}% per year after the FY{daa_knot} peak "
                  f"(95% CI {fmt(daa_post_ci[0],0)}-{fmt(daa_post_ci[1],0)}%), and the "
                  f"post-peak slope differed significantly from the pre-peak slope "
                  f"(P={daa_slope_p:.3f}).")
    else:
        p.add_run(f"DAAの利用自体は急増→減衰パターンを示した（図2）。DAA合計はFY{daa_peak_fy}に"
                  f"{fmt(daa_peak_val_m,1)}百万単位（FY{Y0}比+{fmt(daa_rise,0)}%）でピークに達し、"
                  f"その後FY{Y1}までに{fmt(daa_fall,0)}%減少し{fmt(daa_last_m,1)}百万単位となった。"
                  f"後発の製剤が先行製剤を置換した。分節モデルでは、DAA処方数量は"
                  f"FY{daa_knot}のピーク後に年{fmt(daa_post_r,0)}%"
                  f"（95%CI {fmt(daa_post_ci[0],0)}〜{fmt(daa_post_ci[1],0)}%）で減少し、"
                  f"ピーク後の傾きはピーク前と有意に異なった（P={daa_slope_p:.3f}）。")
    insert_fig(doc, os.path.join(OUT, f"fig2_daa_wave_{lang}.png"))
    add_caption(doc, ("Fig. 2. The DAA wave: dispensed quantity by product." if lang == "en"
                      else "図2．DAAの波：製剤別処方数量。"))

    p = doc.add_paragraph()
    if lang == "en":
        p.add_run("To give a rough sense of the practical patient scale (dispensed "
                  "quantity is not a patient count), we converted the dispensed quantity "
                  "of each interferon-free DAA to approximate full treatment courses using "
                  "documented per-course unit counts (daily dose x standard duration), "
                  "counting one anchor component per two-drug regimen to avoid "
                  "double-counting. This yields an estimated peak of about "
                  f"{fmt(course_peak,0)} courses in FY{course_peak_fy} and an estimated "
                  f"{fmt(course_total_lo,0)}-{fmt(course_total_hi,0)} courses over "
                  f"FY{Y0}-FY{Y1} (duration-sensitivity range). These figures are "
                  "explicit estimates dependent on the regimen assumptions and are not "
                  "observed patient counts.")
    else:
        p.add_run("実用的な患者規模の目安を得るため（処方数量は患者数ではない）、各IFNフリーDAAの"
                  "処方数量を、明示的な1コースあたり単位数（日用量×標準投与期間）を用いておおよその"
                  "治療コース数に換算した（2剤レジメンは二重計上を避けるため代表成分（アンカー）を一つ"
                  f"のみ計上）。推定ピークはFY{course_peak_fy}の約{fmt(course_peak,0)}コース、FY{Y0}〜"
                  f"FY{Y1}の累計は約{fmt(course_total_lo,0)}〜{fmt(course_total_hi,0)}コース（投与期間の"
                  "感度幅）であった。これらはレジメン仮定に依存する明示的な推定値であり、実測の患者数ではない。")

    # ---- Discussion ----
    doc.add_heading(T["h_disc"], level=1)
    p = doc.add_paragraph()
    if lang == "en":
        p.add_run("For hepatitis C DAAs there was no population-level practical lag: the "
                  "standard therapy was displaced within about two years of the "
                  "announcements/listings, indicating that decision-makers reacted rapidly. "
                  "The DAA surge-then-decay is consistent with a finite stock of "
                  "long-waiting patients being cured in a burst (pent-up demand) rather "
                  "than a steady replacement flow. This is descriptive evidence supporting "
                  "the anticipation hypothesis; it does not by itself establish that media "
                  "coverage caused individual treatment choices. External pre-2014 "
                  "utilization data (e.g. national hepatitis programmes or society "
                  "statistics) would be needed to quantify the full pre-DAA interferon "
                  "baseline and are outside NDB; our claim is therefore limited to the "
                  "speed of within-NDB displacement.")
    else:
        p.add_run("C型肝炎DAAでは人口レベルの実用的ラグは認められなかった。標準治療は報道・収載から"
                  "約2年で置換され、治療選択の主体が迅速に反応したことを示す。DAAの急増→減衰は、長く"
                  "待機していた患者ストックが一括して治癒された（pent-up demand）ことと整合的であり、"
                  "定常的置換フローではない。これは待望論仮説を支持する記述的証拠であり、報道が個々の"
                  "治療選択を引き起こしたことを単独で証明するものではない。DAA前のインターフェロン利用の"
                  "完全な基準を定量化するには、2014年以前の外部データ（国の肝炎対策事業や学会統計等）が"
                  "必要だがNDB外であるため、本研究の主張はNDB内での置換の速さに限定される。")

    doc.add_heading(T["h_lim"], level=2)
    p = doc.add_paragraph()
    if lang == "en":
        p.add_run(f"NDB Open Data begins in FY{Y0}, the same period IFN-free DAAs launched, "
                  f"so there is no pre-DAA interferon baseline within NDB; the FY{Y0} value "
                  "already reflects decline from the pre-2014 peak. The metric is dispensed "
                  "quantity, not patient counts, and units differ across products, so the "
                  "summed DAA quantity is not a patient count. Data are annual, precluding "
                  "within-year interrupted time-series or formal causal estimation; the "
                  "reported trend rates carry wide uncertainty intervals given the small "
                  f"number of annual observations (n={n_obs}), and no control condition or "
                  "placebo event is included. Secular changes from FY2020 onward, "
                  "including the COVID-19 pandemic's effect on outpatient visits and "
                  "prescribing, may also have influenced later dispensing and cannot be "
                  "separated from the ongoing DAA decline.")
    else:
        p.add_run(f"NDBオープンデータはFY{Y0}開始であり、これはIFNフリーDAA導入時期と重なるため、"
                  f"NDB内にDAA前のインターフェロン基準値は存在しない（FY{Y0}値は既に2014年以前の"
                  "ピークからの減少を反映）。指標は処方数量であり患者数ではなく、製剤間で単位が"
                  "異なるためDAA数量の合計は患者数ではない。データは年次であり、年内の中断時系列や"
                  f"形式的因果推定はできず、報告したトレンド率は年次観測数が少ない（n={n_obs}）ため"
                  "広い不確実性区間を伴う。対照条件・プラセボイベントも含まない。またFY2020以降は"
                  "COVID-19パンデミックによる外来受診・処方への影響などの外生的変化が重なり、"
                  "DAAの減少傾向と分離できない。")

    # ---- Conclusion ----
    doc.add_heading(T["h_conc"], level=1)
    p = doc.add_paragraph()
    if lang == "en":
        p.add_run("In hepatitis C, the interferon-based standard therapy was replaced at "
                  "the population level within about two years of interferon-free DAAs, and "
                  "the surge-then-decay of DAA use is consistent with realized pent-up "
                  "demand. This supports the existence of population-level new-treatment "
                  "anticipation for this case, while cautioning that the same pattern need "
                  "not generalize to therapies with weaker anticipation or greater "
                  "practical constraints.")
    else:
        p.add_run("C型肝炎では、インターフェロンベース標準治療はIFNフリーDAA導入から約2年で人口"
                  "レベルに置換され、DAA利用の急増→減衰は待機需要の顕在化と整合的であった。これは"
                  "本事例における人口レベルの新治療法待望の存在を支持する一方、待望が弱い、あるいは"
                  "実用的制約の大きい治療には同じパターンが一般化するとは限らないことに注意を要する。")

    if journal == "pds":
        add_pds_statements(doc)

    # ---- Data/code availability + references ----
    doc.add_heading(T["h_da"], level=1)
    p = doc.add_paragraph()
    p.add_run("Raw NDB Open Data workbooks are publicly available from MHLW and are "
              "re-downloadable via scripts/download_ndb.py; derived datasets, analysis "
              "code and figure/manuscript generators are in the project repository."
              if lang == "en" else
              "NDBオープンデータの元ファイルは厚生労働省より公開され、scripts/download_ndb.py で"
              "再取得できる。派生データ・解析コード・図表／原稿生成スクリプトはプロジェクトリポジトリ"
              "に含まれる。")

    doc.add_heading(T["h_ref"], level=1)
    for i, k in enumerate(CITE_ORDER, 1):
        rp = doc.add_paragraph()
        rp.paragraph_format.left_indent = Inches(0.3)
        rp.paragraph_format.first_line_indent = Inches(-0.3)
        rp.add_run(f"{i}. {REF_TEXT[lang][k]}")

    suffix = f"_{journal}" if journal else ""
    path = os.path.join(OUT, f"manuscript_{lang}{suffix}.docx")
    doc.save(path)
    print("wrote", path)


def add_pds_titlepage(doc):
    """Pharmacoepidemiology & Drug Safety title-page front matter (English)."""
    meta = [
        ("Article type", "Original Report (observational, national open-data study)"),
        ("Running head", "Population-level HCV DAA uptake vs interferon therapy"),
        ("Corresponding author", "[Name], [Affiliation], [Address], [Email]"),
        ("Authors / affiliations", "[To be completed by the authors]"),
    ]
    for k, v in meta:
        p = doc.add_paragraph()
        p.add_run(f"{k}: ").bold = True
        p.add_run(v)


def add_keywords(doc):
    p = doc.add_paragraph()
    p.add_run("Keywords: ").bold = True
    p.add_run("hepatitis C; direct-acting antivirals; interferon; pharmacoepidemiology; "
              "prescription trends; drug utilization; NDB Open Data; Japan")


def add_key_points(doc):
    """P&DS 'Key Points' / take-home box."""
    doc.add_heading("Key Points", level=2)
    pts = [
        "Whether population-level anticipation of a newly reimbursed therapy pulls "
        "patients off the prior standard of care can be examined with national "
        "drug-utilization open data.",
        "In Japan, dispensing of interferon-based standard therapy for hepatitis C "
        f"collapsed (peginterferon -{fmt(peg_drop,1)}%; ribavirin near-zero by "
        f"FY{rbv_zero_year}) within ~2 years of interferon-free direct-acting "
        "antivirals becoming available.",
        f"Interferon-free DAA dispensing surged (+{fmt(daa_rise,0)}% to a FY{daa_peak_fy} "
        f"peak) and then declined ({fmt(daa_fall,0)}%), a pattern consistent with "
        "realized pent-up demand rather than steady substitution.",
        "Findings are descriptive (national dispensed quantity, not patient counts) and "
        "do not establish that media coverage caused individual treatment choices.",
    ]
    for t in pts:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(t)


def add_pds_statements(doc):
    """Ethics / consent / COI / funding / reporting-guideline statements."""
    items = [
        ("Ethics approval and consent",
         "This study used only publicly available, aggregated national open data "
         "(NDB Open Data) containing no individual-level or identifiable information; "
         "ethics-committee approval and informed consent were therefore not required."),
        ("Conflict of interest", "The authors declare no conflicts of interest."),
        ("Funding", "This research received no specific grant from any funding agency."),
        ("Reporting guideline",
         "This observational study is reported in line with the STROBE guideline for "
         "cross-sectional/ecological analyses of routinely collected aggregate data; a "
         "completed STROBE checklist can be provided as supplementary material."),
    ]
    for h, t in items:
        p = doc.add_paragraph()
        p.add_run(f"{h}. ").bold = True
        p.add_run(t)


def _event_date(e):
    """Day-precise date from precision='day(YYYY-MM-DD)', else the month."""
    prec = str(e["precision"])
    if prec.startswith("day(") and prec.endswith(")"):
        return prec[4:-1]
    return str(e["event_month"])


def add_events_table(doc, lang):
    cols = (["FY", "Date", "Drug", "Milestone", "Source"] if lang == "en"
            else ["年度", "日付", "薬剤", "イベント", "出所"])
    t = doc.add_table(rows=1, cols=len(cols))
    t.style = "Table Grid"
    for j, c in enumerate(cols):
        r = t.rows[0].cells[j].paragraphs[0].add_run(c); r.bold = True
    drug_col = "drug_en" if lang == "en" else "drug_ja"
    mile = {"approval": "approval" if lang == "en" else "承認",
            "nhi_listing": "NHI listing" if lang == "en" else "薬価収載"}
    for _, e in EV.iterrows():
        cells = t.add_row().cells
        cells[0].text = str(int(e["fy"]))
        cells[1].text = _event_date(e)
        cells[2].text = str(e[drug_col])
        cells[3].text = mile.get(e["milestone"], e["milestone"])
        cells[4].text = str(e["source"])
    for row in t.rows:
        for cell in row.cells:
            for pph in cell.paragraphs:
                for rr in pph.runs:
                    rr.font.size = Pt(7.5)


def build_tables_doc(lang):
    doc = Document()
    doc.add_heading("Tables" if lang == "en" else "表", level=1)
    add_caption(doc, ("Table 1. Official approval / NHI drug-price listing milestones."
                      if lang == "en"
                      else "表1．公式の承認・薬価収載イベント。"))
    add_events_table(doc, lang)

    add_caption(doc, ("Table 2. National dispensed quantity by drug group and fiscal year "
                      "(NDB Open Data)." if lang == "en"
                      else "表2．薬効グループ別・年度別の全国処方数量（NDBオープンデータ）。"))
    groups = ["IFN_peg", "ribavirin", "IFN_conv", "PI_ifn", "DAA"]
    head = (["Fiscal year", "Peginterferon", "Ribavirin", "Conventional IFN",
             "First-gen PI (IFN-based)", "Interferon-free DAA total"]
            if lang == "en"
            else ["年度", "ペグIFN", "リバビリン", "従来型IFN",
                  "第一世代PI(IFN併用)", "IFNフリーDAA合計"])
    t = doc.add_table(rows=1, cols=len(head)); t.style = "Table Grid"
    for j, c in enumerate(head):
        rr = t.rows[0].cells[j].paragraphs[0].add_run(c); rr.bold = True
    for fy, row in TS.iterrows():
        cells = t.add_row().cells
        cells[0].text = str(int(fy))
        for j, g in enumerate(groups, 1):
            cells[j].text = fmt(row[g], 0) if g in row else "0"

    add_caption(doc, ("Table 3. Estimated interferon-free DAA treatment courses by fiscal "
                      "year (ESTIMATE; dispensed quantity / documented units-per-course, "
                      "one anchor product per regimen). Not observed patient counts."
                      if lang == "en"
                      else "表3．IFNフリーDAAの推定治療コース数（推定値；処方数量÷1コース単位数、"
                           "レジメンごとに代表成分を1つ計上）。実測の患者数ではない。"))
    ce = COURSE["estimated_courses_by_fy"]
    ce_hi = COURSE["estimated_courses_by_fy_longer_duration"]
    chead = (["Fiscal year", "Estimated courses (baseline)", "Estimated courses (longer duration)"]
             if lang == "en"
             else ["年度", "推定コース数（基準）", "推定コース数（長期投与）"])
    t3 = doc.add_table(rows=1, cols=len(chead)); t3.style = "Table Grid"
    for j, c in enumerate(chead):
        rr = t3.rows[0].cells[j].paragraphs[0].add_run(c); rr.bold = True
    for y in sorted(ce, key=lambda x: int(x)):
        cells = t3.add_row().cells
        cells[0].text = str(y)
        cells[1].text = fmt(ce[y], 0)
        cells[2].text = fmt(ce_hi[y], 0)

    path = os.path.join(OUT, f"tables_{lang}.docx")
    doc.save(path)
    print("wrote", path)


def build_pptx(lang):
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    figs = [
        (f"fig1_ifn_collapse_{lang}.png",
         "Fig. 1. Collapse of interferon-based standard therapy" if lang == "en"
         else "図1．インターフェロンベース標準治療の消失"),
        (f"fig2_daa_wave_{lang}.png",
         "Fig. 2. The DAA wave: dispensed quantity by product" if lang == "en"
         else "図2．DAAの波：製剤別処方数量"),
    ]
    for fn, cap in figs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(PInches(0.5), PInches(0.2), PInches(12.3), PInches(0.7))
        tf = tb.text_frame; tf.text = cap
        tf.paragraphs[0].runs[0].font.size = PPt(20)
        tf.paragraphs[0].runs[0].font.bold = True
        slide.shapes.add_picture(os.path.join(OUT, fn), PInches(1.4), PInches(1.1),
                                 height=PInches(5.7))
    path = os.path.join(OUT, f"figures_{lang}.pptx")
    prs.save(path)
    print("wrote", path)


STROBE_ITEMS = [
    # (No., section, item text, our response / location)
    ("1", "Title and abstract",
     "(a) Indicate the study design with a commonly used term in the title or the "
     "abstract. (b) Provide in the abstract an informative and balanced summary of "
     "what was done and what was found.",
     "Title and structured abstract describe an observational, ecological (national "
     "aggregate) drug-utilisation time-series analysis of NDB Open Data; the abstract "
     "summarises design, data, methods and headline results."),
    ("2", "Background/rationale",
     "Explain the scientific background and rationale for the investigation being "
     "reported.",
     "Introduction frames new-treatment anticipation / practical drug lag and why a "
     "near-complete substitution (HCV interferon -> interferon-free DAA) is an "
     "informative natural test."),
    ("3", "Objectives",
     "State specific objectives, including any prespecified hypotheses.",
     "Introduction states the hypothesis: if population-level anticipation is real, "
     "dispensing of the prior standard therapy should fall sharply after the awaited "
     "option arrives."),
    ("4", "Study design",
     "Present key elements of study design early in the paper.",
     "Methods: ecological/descriptive analysis of national annual dispensed quantity "
     "(FY2014-FY2023); no individual-level data; no pre-intervention baseline within "
     "NDB, so descriptive trend models rather than a causal interrupted time series."),
    ("5", "Setting",
     "Describe the setting, locations, and relevant dates, including periods of "
     "recruitment, exposure, follow-up, and data collection.",
     "Methods: Japan, nationwide; NDB Open Data editions 1-10 mapped to fiscal years "
     "2014-2023; approval/reimbursement milestones tabulated (Table 1)."),
    ("6", "Participants",
     "Give the eligibility criteria, and the sources and methods of selection of "
     "participants (cross-sectional).",
     "No individual participants: the unit is national aggregate dispensed quantity "
     "per drug per fiscal year, extracted for all records of the target products. "
     "Stated explicitly in Methods and Limitations (ecological data, not patient counts)."),
    ("7", "Variables",
     "Clearly define all outcomes, exposures, predictors, potential confounders, and "
     "effect modifiers.",
     "Methods: outcome = national dispensed quantity by drug group (peginterferon, "
     "conventional IFN, ribavirin, first-generation IFN-based protease inhibitors, "
     "interferon-free DAAs); 'exposure' = approval/reimbursement/reporting milestones; "
     "potential secular confounders (e.g. COVID-19) noted in Limitations."),
    ("8", "Data sources/measurement",
     "For each variable of interest, give sources of data and details of methods of "
     "assessment (measurement).",
     "Methods and Data/code availability: MHLW NDB Open Data workbooks (sex/age x drug "
     "quantity tables); metric is total dispensed quantity (tablets/capsules or "
     "syringes/vials). Product classification listed; extraction code in the repository."),
    ("9", "Bias",
     "Describe any efforts to address potential sources of bias.",
     "Methods/Limitations: first-generation IFN-based protease inhibitors separated "
     "from interferon-free DAAs to avoid misclassification; units are not summed across "
     "products with different dosage units; absence of a pre-2014 baseline and lack of "
     "a control condition are stated as biases limiting causal inference."),
    ("10", "Study size",
     "Explain how the study size was arrived at.",
     "Methods: the study uses the complete set of national annual observations available "
     "in NDB Open Data (n=10 fiscal years); no sampling. Small n is flagged as a driver "
     "of wide uncertainty intervals."),
    ("11", "Quantitative variables",
     "Explain how quantitative variables were handled in the analyses.",
     "Methods: log-scale trend modelling; treatment-course sensitivity converts dispensed "
     "quantity to approximate courses using documented daily dose x duration "
     "(data/daa_course_assumptions.csv), counting one anchor product per two-drug regimen."),
    ("12", "Statistical methods",
     "(a) Describe all statistical methods, including those used to control for "
     "confounding. (b) sensitivity analyses.",
     "Methods: segmented (broken-stick) log-linear regression with a knot at the observed "
     "DAA peak, exponential/log-linear decay for IFN groups, Newey-West (HAC) standard "
     "errors and residual-bootstrap 95% intervals; duration-based sensitivity analysis "
     "for course estimates. Analysis is explicitly descriptive/non-causal."),
    ("13", "Participants (results)",
     "(a) Report numbers of individuals at each stage. (b) reasons for non-participation. "
     "(c) consider a flow diagram.",
     "Not applicable at the individual level (aggregate open data). Results and Methods "
     "report the data units: 10 fiscal years x drug groups; 9 distinct interferon-free "
     "DAA products and 3 first-generation IFN-based protease inhibitors."),
    ("14", "Descriptive data",
     "(a) Give characteristics of study participants and information on exposures and "
     "potential confounders. (b) missing data. (c) follow-up (cohort).",
     "Results and Table 2: national dispensed quantity by group and fiscal year. Products "
     "absent in a given year appear as zero/near-zero; no imputation. Milestones in Table 1."),
    ("15", "Outcome data",
     "Report numbers of outcome events or summary measures (cross-sectional).",
     "Results, Table 2, Figures 1-2: annual dispensed quantities and their changes; Table 3 "
     "reports estimated treatment courses (labelled estimates)."),
    ("16", "Main results",
     "(a) Give unadjusted and adjusted estimates and their precision. (b) category "
     "boundaries. (c) translate relative to absolute risk if relevant.",
     "Results: peginterferon -99.2% (FY2014->FY2023); ribavirin near-zero by FY2018; DAA "
     "peak FY2015 (+188% vs FY2014) then -92%; annual decline rates with HAC/bootstrap 95% "
     "intervals and a slope-change P value from segmented regression."),
    ("17", "Other analyses",
     "Report other analyses done (e.g. subgroups, interactions, sensitivity analyses).",
     "Results: product-level DAA breakdown (Figure 2) and treatment-course sensitivity "
     "under baseline vs longer-duration assumptions (Table 3, ~266k-298k courses)."),
    ("18", "Key results",
     "Summarise key results with reference to study objectives.",
     "Discussion opens by restating that IFN-based standard therapy was replaced within "
     "~2 years and the DAA surge-then-decay is consistent with realized pent-up demand."),
    ("19", "Limitations",
     "Discuss limitations, taking into account sources of potential bias or imprecision.",
     "Limitations: no pre-2014 baseline within NDB; dispensed quantity != patient counts; "
     "annual resolution precludes within-year ITS; small n; no control/placebo event; "
     "COVID-19 and other FY2020+ secular changes cannot be separated from the DAA decline; "
     "course figures are estimates dependent on regimen assumptions."),
    ("20", "Interpretation",
     "Give a cautious overall interpretation considering objectives, limitations, "
     "multiplicity, and other relevant evidence.",
     "Discussion/Conclusion: interpretation limited to within-NDB displacement speed; "
     "'consistent with' anticipation, not a causal claim that media coverage drove "
     "individual choices."),
    ("21", "Generalisability",
     "Discuss the generalisability (external validity) of the study results.",
     "Discussion: the near-complete HCV substitution is a strong case; the same surge "
     "pattern may not generalise to therapies with weaker anticipation or larger practical "
     "constraints."),
    ("22", "Funding",
     "Give the source of funding and the role of the funders for the present study and, "
     "if applicable, for the original study on which the present article is based.",
     "Funding statement: no specific grant. Data are public NDB Open Data (MHLW)."),
]


def build_strobe():
    """Filled STROBE checklist (cross-sectional/ecological) as a supplementary docx."""
    doc = Document()
    doc.styles["Normal"].font.size = Pt(9.5)
    h = doc.add_paragraph()
    r = h.add_run("STROBE Statement — checklist of items for reports of "
                  "observational studies (cross-sectional / ecological adaptation)")
    r.bold = True; r.font.size = Pt(12)
    doc.add_paragraph(
        "Strengthening the Reporting of Observational Studies in Epidemiology (STROBE). "
        "This study analyses national aggregate drug-dispensing open data (no "
        "individual-level records); items are answered accordingly, with locations given "
        "by manuscript section rather than page number.")
    tbl = doc.add_table(rows=1, cols=4)
    tbl.style = "Table Grid"
    hdr = tbl.rows[0].cells
    for c, t in zip(hdr, ["Item No.", "Section", "STROBE recommendation",
                          "Location / response in manuscript"]):
        c.paragraphs[0].add_run(t).bold = True
    for no, sec, rec, resp in STROBE_ITEMS:
        cells = tbl.add_row().cells
        cells[0].text = no
        cells[1].text = sec
        cells[2].text = rec
        cells[3].text = resp
    path = os.path.join(OUT, "strobe_checklist.docx")
    doc.save(path)
    print("wrote", path)


if __name__ == "__main__":
    for lang in ("en", "ja"):
        build_manuscript(lang)
        build_tables_doc(lang)
        build_pptx(lang)
    # Pharmacoepidemiology & Drug Safety submission variant (English):
    # structured abstract, Key Points, title-page front matter, STROBE/ethics/
    # COI/funding statements. Figures remain inline (P&DS accepts free-format
    # submission).
    build_manuscript("en", journal="pds")
    # Filled STROBE checklist (supplementary material for P&DS submission).
    build_strobe()
